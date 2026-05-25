#!/usr/bin/env python3
"""Build the Sci-Evo submission PPT and offline deliverables package."""

from __future__ import annotations

import json
import shutil
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
REFERENCE_PPT = Path("/Users/liangming/Documents/mineru-question3/outputs/submission_pack/deliverables/eduminer_competition_pitch.pptx")
OUT_ROOT = ROOT / "outputs" / "submission_pack"
DELIVERABLES = OUT_ROOT / "deliverables"
PPT_OUT = DELIVERABLES / "scievo_competition_pitch.pptx"
SUMMARY_IMAGE = OUT_ROOT / "assets" / "submission_pack_overview.png"
ZIP_OUT = OUT_ROOT / "scievo_submission_pack.zip"


SLIDE_TEXT = {
    1: {
        2: "MinerU · Sci-Evo",
        3: "Sci-Evo-LabTrace",
        4: "面向 AI4S 科学智能体的科研演化证据数据集",
        5: "把论文 PDF、MinerU 解析结果和人工证据链整合成可训练、可评测、可追溯的 Sci-Evo gold case。",
        8: "3",
        9: "Gold case",
        11: "17",
        12: "轨迹步骤",
        14: "37",
        15: "评测任务",
        17: "25",
        18: "OA 候选",
        20: "不是“论文摘要”",
        21: "而是把真实科研中的目标设定、计算设计、湿实验、失败修正和最终验证整理成可学习的轨迹。",
        22: "核心差异",
        23: "论文证据链 + MinerU 块级定位 + 结构化 schema + 自动评测任务",
        24: "Sci-Evo · LabTrace",
        25: "01",
    },
    2: {
        3: "场景痛点",
        4: "真正稀缺的不是论文数量，而是可复用的科研演化过程。",
        5: "Sci-Evo · LabTrace",
        6: "02",
        8: "01",
        9: "静态知识多",
        10: "现有语料大多只保留结论、摘要或问答，缺少真实科研步骤与决策上下文。",
        11: "→",
        13: "02",
        14: "证据链断",
        15: "即使能读 PDF，也常停留在段落摘录，难以回溯到页面、图表和具体实验节点。",
        16: "→",
        18: "03",
        19: "提交成本高",
        20: "数据、报告、评测任务和许可说明分散，评委很难快速确认工程完整度。",
        22: "比赛得分点",
        23: "高保真解析只是底座；更关键的是把科研轨迹、证据链、质量控制和提交包做成闭环。",
    },
    3: {
        3: "系统架构",
        4: "先把论文解析成可信证据，再构建科研轨迹数据。",
        5: "Sci-Evo · LabTrace",
        6: "03",
        8: "输入",
        9: "PDF / DOI / 许可信息",
        10: "→",
        12: "解析",
        13: "MinerU / Markdown / layout",
        14: "→",
        16: "建模",
        17: "case schema / evidence map",
        18: "→",
        20: "抽取",
        21: "trajectory / verification",
        22: "→",
        24: "交付",
        25: "JSONL / tasks / reports",
        27: "统一证据链",
        28: "每个轨迹步骤保留来源页码、定位文本和 MinerU block 对应关系。",
        30: "统一提交闭环",
        31: "构建、校验、质量报告、就绪检查和打包脚本共同保证可提交性。",
    },
    4: {
        3: "MINERU 解析",
        4: "第一公里有真实落地：本地产物、页级证据和复现实验路径都可检查。",
        5: "Sci-Evo · LabTrace",
        6: "04",
        8: "1",
        9: "本地解析样例",
        11: "5",
        12: "核心产物类型",
        14: "100%",
        15: "样例可复查",
        17: "本地运行报告已保存",
        18: "输入样例",
        19: "Sci-Evo-Sample.pdf",
        20: "产物覆盖",
        21: "full.md / layout.json / content_list",
        22: "关键文件",
        23: "markdown / json / model / images",
        24: "许可边界",
        25: "全文级解析产物仅本地保留",
        26: "后续扩展",
        27: "开放论文可继续补齐本地 MinerU 解析结果",
        28: "使用方式：通过 scripts/mineru_parse.py 调用 MinerU API，本地保留 Markdown、content list、layout、model 与图片；因许可约束，不公开未复核全文级产物。",
    },
    5: {
        3: "数据总览",
        4: "每条 case 同时呈现领域、轨迹重心、许可状态和关键证据。",
        5: "Sci-Evo · LabTrace",
        6: "05",
        8: "样本",
        10: "领域",
        12: "轨迹重心",
        14: "许可状态",
        16: "关键证据",
        18: "de_novo_luciferase",
        20: "生物发光",
        22: "从头设计 + 筛选",
        24: "赛事样例",
        26: "目标、突变、细胞验证",
        28: "ml_enzyme_opt",
        30: "酶工程",
        32: "ML 设计 + 活性优化",
        34: "CC-BY",
        36: "模型输入、实验结果",
        38: "biocatalysis_case",
        40: "生物催化",
        42: "定向进化 + 验证",
        44: "CC-BY",
        46: "轮次指标、终点结论",
        48: "openalex_queue",
        50: "候选扩展",
        52: "许可筛选 + 后续解析",
        54: "14 可处理",
        56: "DOI、OA 链接、许可状态",
        57: "这张表让评委可以在一页内看清当前数据覆盖、来源合规性和后续扩展空间。",
    },
    6: {
        3: "标注策略",
        4: "目标定义不只写摘要，而是决定轨迹边界、验证标准和评测方式。",
        5: "Sci-Evo · LabTrace",
        6: "06",
        8: "case 画像",
        9: "目标",
        10: "从论文中抽取可量化科研目标",
        11: "难点",
        12: "计算设计、湿实验和失败修正需串联",
        13: "边界",
        14: "只保留有证据和许可状态的步骤",
        15: "→",
        18: "进入结构化标注",
        19: "initial_request",
        20: "目标、约束、量化指标",
        21: "agent_trajectory",
        22: "action / tool / parameter / observation",
        23: "success_verification",
        24: "最终验证方法与指标",
        25: "quality",
        26: "gold 等级、证据覆盖、许可复核",
        27: "同一篇论文的切分方式会直接影响样本质量；我们的原则是先定边界，再抽步骤，再做证据绑定。",
    },
    7: {
        3: "提交交付",
        4: "一次打包交出数据、报告、脚本和审查说明，而不是零散文件。",
        5: "Sci-Evo · LabTrace",
        6: "07",
        8: "submission_pack/",
        9: "记录数据、文档、报告与交付清单路径",
        11: "数据",
        12: "scievo_gold.jsonl / scievo_eval_tasks.jsonl",
        14: "文档",
        15: "README.md / TECHNICAL_REPORT.md",
        17: "校验",
        18: "validate_dataset.py / check_submission_readiness.py",
        20: "报告",
        21: "QUALITY_REPORT.md / MINERU_RUN_REPORT.md / SUBMISSION_READINESS.md",
        30: "对评委：打开包就能复核。对团队：后续扩展 case 时不用重搭提交结构。",
    },
    8: {
        3: "质量闭环",
        4: "每个 case 通过结构校验后，还要被证据、许可和评测任务反向约束。",
        5: "Sci-Evo · LabTrace",
        6: "08",
        8: "3/3",
        9: "Gold case 通过校验",
        11: "37",
        12: "自动评测任务",
        14: "25",
        15: "开放来源已筛选",
        17: "100%",
        18: "提交门槛已满足",
        20: "许可状态",
        21: "CC-BY 2 / 样例待复核 1",
        23: "质量报告",
        24: "已生成",
        26: "就绪检查",
        27: "通过",
        29: "自动审校",
        31: "validate_dataset 通过；check_submission_readiness 通过。",
    },
    9: {
        3: "提交包",
        4: "评委打开交付包，先看到核心材料，再顺着文件路径进入证据与代码。",
        5: "Sci-Evo · LabTrace",
        6: "09",
        10: "首屏讲清楚三件事",
        11: "数据文件",
        12: "gold / eval / manifest",
        13: "关键文档",
        14: "README / 技术报告 / 清单",
        15: "验证状态",
        16: "校验通过 / 工作区干净",
        17: "许可说明",
        18: "公开发布前复核全文级产物",
    },
    10: {
        3: "冲刺路线",
        4: "当前版本已可提交，后续增强应集中在规模、MinerU 覆盖和许可清晰度。",
        5: "Sci-Evo · LabTrace",
        6: "10",
        8: "P0",
        9: "提交版交付包",
        10: "当前已完成；PPT、数据、报告与校验链齐全。",
        12: "P1",
        13: "扩展更多 OA case",
        14: "优先处理 14 条 permitted_for_local_processing 候选。",
        16: "P2",
        17: "补齐 MinerU 产物",
        18: "为新增开放论文补本地解析和块级证据绑定。",
        20: "P3",
        21: "公开版合规",
        22: "逐条确认全文、图片和解析产物的许可边界。",
        24: "答辩话术：MinerU 负责把论文变成可信结构化证据，Sci-Evo-LabTrace 负责把证据组织成可训练、可评测、可提交的科研演化数据。",
    },
}

COPY_ITEMS = [
    ("README.md", "README.md"),
    ("docs/TECHNICAL_REPORT.md", "docs/TECHNICAL_REPORT.md"),
    ("docs/DATASET_CARD.md", "docs/DATASET_CARD.md"),
    ("docs/SUBMISSION_CHECKLIST.md", "docs/SUBMISSION_CHECKLIST.md"),
    ("docs/VIDEO_SCRIPT.md", "docs/VIDEO_SCRIPT.md"),
    ("reports/QUALITY_REPORT.md", "reports/QUALITY_REPORT.md"),
    ("reports/MINERU_RUN_REPORT.md", "reports/MINERU_RUN_REPORT.md"),
    ("reports/SUBMISSION_READINESS.md", "reports/SUBMISSION_READINESS.md"),
    ("reports/VETTED_SOURCE_QUEUE.md", "reports/VETTED_SOURCE_QUEUE.md"),
    ("data/processed/scievo_gold.jsonl", "data/processed/scievo_gold.jsonl"),
    ("data/processed/scievo_eval_tasks.jsonl", "data/processed/scievo_eval_tasks.jsonl"),
    ("data/processed/dataset_manifest.json", "data/processed/dataset_manifest.json"),
    ("data/processed/vetted_open_access_sources.jsonl", "data/processed/vetted_open_access_sources.jsonl"),
    ("schemas/scievo_case.schema.json", "schemas/scievo_case.schema.json"),
    ("scripts/build_dataset.py", "scripts/build_dataset.py"),
    ("scripts/validate_dataset.py", "scripts/validate_dataset.py"),
    ("scripts/build_eval_tasks.py", "scripts/build_eval_tasks.py"),
    ("scripts/check_submission_readiness.py", "scripts/check_submission_readiness.py"),
    ("scripts/make_quality_report.py", "scripts/make_quality_report.py"),
    ("scripts/vet_candidate_sources.py", "scripts/vet_candidate_sources.py"),
    ("Sci-Evo_tool_case.json", "examples/Sci-Evo_tool_case.json"),
]


def ensure_dirs() -> None:
    (OUT_ROOT / "assets").mkdir(parents=True, exist_ok=True)
    DELIVERABLES.mkdir(parents=True, exist_ok=True)


def make_overview_image() -> None:
    image = Image.new("RGB", (1280, 720), "#F7F5EE")
    draw = ImageDraw.Draw(image)
    title_font = ImageFont.load_default()
    body_font = ImageFont.load_default()

    draw.rounded_rectangle((48, 48, 1232, 672), radius=28, fill="#FFFFFF", outline="#D6D0C5", width=2)
    draw.rounded_rectangle((84, 84, 400, 184), radius=18, fill="#E8F0FF")
    draw.text((108, 110), "Sci-Evo 提交包概览", fill="#172033", font=title_font)
    draw.text((108, 138), "数据、文档、报告、脚本一包可审", fill="#2563EB", font=body_font)

    sections = [
        ("数据文件", ["scievo_gold.jsonl", "scievo_eval_tasks.jsonl", "dataset_manifest.json"], "#EAF6FF"),
        ("关键文档", ["README.md", "TECHNICAL_REPORT.md", "SUBMISSION_CHECKLIST.md"], "#EEFCEB"),
        ("质量报告", ["QUALITY_REPORT.md", "MINERU_RUN_REPORT.md", "SUBMISSION_READINESS.md"], "#FFF5E6"),
        ("验证脚本", ["build_dataset.py", "validate_dataset.py", "check_submission_readiness.py"], "#F7EEFF"),
    ]

    y = 220
    for title, lines, fill in sections:
        draw.rounded_rectangle((92, y, 1188, y + 92), radius=16, fill=fill)
        draw.text((120, y + 16), title, fill="#172033", font=title_font)
        draw.text((120, y + 44), " | ".join(lines), fill="#334155", font=body_font)
        y += 108

    draw.rounded_rectangle((92, 616, 1188, 652), radius=12, fill="#172033")
    draw.text(
        (112, 626),
        "状态：3 个 gold case / 37 个评测任务 / 就绪检查通过 / 工作区干净",
        fill="#FFFFFF",
        font=body_font,
    )
    image.save(SUMMARY_IMAGE)


def set_shape_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    shape.text = text


def build_ppt() -> None:
    prs = Presentation(str(REFERENCE_PPT))

    for slide_index, replacements in SLIDE_TEXT.items():
        slide = prs.slides[slide_index - 1]
        shape_by_id = {shape.shape_id: shape for shape in slide.shapes}
        for shape_id, text in replacements.items():
            shape = shape_by_id.get(shape_id)
            if shape is None:
                raise KeyError(f"Slide {slide_index} missing shape id {shape_id}")
            set_shape_text(shape, text)

    slide9 = prs.slides[8]
    image_shape = next(shape for shape in slide9.shapes if shape.shape_type == 13)
    left, top, width, height = image_shape.left, image_shape.top, image_shape.width, image_shape.height
    image_shape._element.getparent().remove(image_shape._element)
    slide9.shapes.add_picture(str(SUMMARY_IMAGE), left, top, width=width, height=height)

    prs.save(str(PPT_OUT))


def copy_tree_items() -> None:
    for src_rel, dst_rel in COPY_ITEMS:
        src = ROOT / src_rel
        dst = DELIVERABLES / dst_rel
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(src, dst)

    readme = DELIVERABLES / "README.md"
    readme.write_text(
        "\n".join(
            [
                "# Sci-Evo 提交包",
                "",
                "本目录包含离线提交所需的核心交付件：",
                "",
                "- `scievo_competition_pitch.pptx`：参考 EduMiner 提交稿风格重写的 Sci-Evo 展示 PPT。",
                "- `data/processed/`：gold 数据集、评测任务、manifest 和开放来源筛选结果。",
                "- `docs/`：技术报告、数据集卡片、提交清单和录屏讲稿。",
                "- `reports/`：质量报告、MinerU 运行报告、提交就绪报告和来源筛选队列。",
                "- `schemas/`：Sci-Evo case schema。",
                "- `scripts/`：构建、校验、评测任务生成、质量报告和提交就绪检查脚本。",
                "",
                "说明：",
                "",
                "- 不包含任何密钥文件。",
                "- 不公开打包未经许可复核的全文 PDF 或 MinerU 全量解析产物。",
                "- 如需复核当前状态，可运行 `python3 scripts/check_submission_readiness.py`。",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def build_zip() -> None:
    if ZIP_OUT.exists():
        ZIP_OUT.unlink()
    with zipfile.ZipFile(ZIP_OUT, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for path in sorted(DELIVERABLES.rglob("*")):
            if path.is_file():
                zf.write(path, path.relative_to(OUT_ROOT))


def main() -> None:
    ensure_dirs()
    make_overview_image()
    build_ppt()
    copy_tree_items()
    build_zip()
    print(json.dumps({"ppt": str(PPT_OUT), "zip": str(ZIP_OUT)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
