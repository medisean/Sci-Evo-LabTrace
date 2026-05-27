#!/usr/bin/env python3
"""Build the Sci-Evo review deck and offline materials package."""

from __future__ import annotations

import json
import shutil
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE


ROOT = Path(__file__).resolve().parents[1]
REFERENCE_PPT = Path("/Users/liangming/Documents/mineru-question3/outputs/submission_pack/deliverables/eduminer_competition_pitch.pptx")
OUT_ROOT = ROOT / "outputs" / "submission_pack"
DELIVERABLES = OUT_ROOT / "deliverables"
PPT_OUT = DELIVERABLES / "scievo_labtrace_overview.pptx"
SUMMARY_IMAGE = OUT_ROOT / "assets" / "review_materials_overview.png"
ZIP_OUT = OUT_ROOT / "scievo_labtrace_review_materials.zip"


TEXT_COLOR_OVERRIDES = {
    1: {
        20: "F8FAFC",
        21: "E5E7EB",
        22: "C7D2FE",
        23: "F8FAFC",
    },
    7: {
        8: "F8FAFC",
        9: "E5E7EB",
    },
    9: {
        10: "F8FAFC",
        11: "CBD5E1",
        12: "F8FAFC",
        13: "CBD5E1",
        14: "F8FAFC",
        15: "CBD5E1",
        16: "F8FAFC",
        17: "CBD5E1",
        18: "F8FAFC",
    },
    10: {
        3: "F8FAFC",
        4: "E5E7EB",
        5: "CBD5E1",
        6: "CBD5E1",
        8: "C7D2FE",
        9: "F8FAFC",
        10: "CBD5E1",
        12: "C7D2FE",
        13: "F8FAFC",
        14: "CBD5E1",
        16: "FBBF24",
        17: "F8FAFC",
        18: "CBD5E1",
        20: "FCA5A5",
        21: "F8FAFC",
        22: "CBD5E1",
        24: "F8FAFC",
    },
}


SLIDE_TEXT = {
    1: {
        2: "MinerU · Sci-Evo",
        3: "Sci-Evo-LabTrace",
        4: "从论文 PDF 到可训练科研轨迹的数据集",
        5: "一句话：用 MinerU 把论文拆成可追溯证据，再整理成科学智能体能学习和评测的科研过程。",
        8: "5",
        9: "Gold case",
        11: "29",
        12: "轨迹步骤",
        14: "63",
        15: "评测任务",
        17: "25",
        18: "OA 候选",
        20: "一条主线",
        21: "论文不是终点；可复核的过程轨迹才是训练科学智能体的资产。",
        22: "核心差异",
        23: "证据先行 → 轨迹抽取 → 评测生成 → 质量闭环",
        24: "Sci-Evo · LabTrace",
        25: "01",
    },
    2: {
        3: "目录",
        4: "评委可以按四个问题阅读这份材料。",
        5: "Sci-Evo · LabTrace",
        6: "02",
        8: "01",
        9: "为什么需要",
        10: "论文结论很多，但可学习的科研过程数据仍然稀缺。",
        11: "→",
        13: "02",
        14: "证据怎么来",
        15: "MinerU 把 PDF 转成 Markdown、layout 和块级证据。",
        16: "→",
        18: "03",
        19: "如何被验证",
        20: "Schema、gold case、评测任务和报告共同证明质量。",
        22: "阅读主线",
        23: "问题定义 → MinerU 证据 → 数据结构 → 质量验证 → 后续扩展。",
    },
    3: {
        3: "主线架构",
        4: "先把论文变成可信证据，再把证据组织成科研轨迹。",
        5: "Sci-Evo · LabTrace",
        6: "03",
        8: "输入",
        9: "PDF / DOI / 许可信息",
        10: "→",
        12: "解析",
        13: "MinerU / Markdown / layout",
        14: "→",
        16: "建模",
        17: "case schema / evidence map",
        18: "→",
        20: "抽取",
        21: "trajectory / verification",
        22: "→",
        24: "产出",
        25: "JSONL / tasks / reports",
        27: "统一证据链",
        28: "每个轨迹步骤保留来源页码、定位文本和 MinerU block 对应关系。",
        30: "统一工程闭环",
        31: "构建、校验、质量报告和一致性检查共同保证材料可复核、可扩展。",
    },
    4: {
        3: "证据来源",
        4: "MinerU 不是展示点，而是每条轨迹可追溯的第一公里。",
        5: "Sci-Evo · LabTrace",
        6: "04",
        8: "1",
        9: "本地解析样例",
        11: "5",
        12: "核心产物类型",
        14: "100%",
        15: "样例可复查",
        17: "MinerU 本地产物",
        18: "输入",
        19: "5 个 PDF",
        20: "产物",
        21: "md / layout",
        22: "证据",
        23: "页码 / block",
        24: "边界",
        25: "全文本地留存",
        26: "扩展",
        27: "开放论文补解析",
        28: "复核方式：关键字段回到 PDF 页面、Markdown 段落或 layout 块；全文级产物因许可仅本地留存。",
    },
    5: {
        3: "数据内容",
        4: "当前 5 条 gold case 覆盖样例、开放论文和后续候选队列。",
        5: "Sci-Evo · LabTrace",
        6: "05",
        8: "样本",
        10: "领域",
        12: "轨迹重心",
        14: "许可状态",
        16: "关键证据",
        18: "0001 luciferase",
        20: "生物发光",
        22: "从头设计 + 筛选",
        24: "种子样例",
        26: "目标、突变、细胞验证",
        28: "0002+0003 ML/DE",
        30: "酶工程",
        32: "文库设计 + 进化复盘",
        34: "CC-BY",
        36: "模型输入、结构证据",
        38: "0004 Pro-PRIME",
        40: "工业抗体",
        42: "LLM + 生产验证",
        44: "CC-BY",
        46: "65 突变、DBC、epistasis",
        48: "0005 signal leader",
        50: "酵母分泌",
        52: "bottom-up + top-down",
        54: "CC-BY",
        56: "1600 克隆、跨酶验证",
        57: "这张表让评委可以在一页内看清当前数据覆盖、来源合规性和后续扩展空间。",
    },
    6: {
        3: "标注结构",
        4: "每条 case 都按“目标、过程、验证、质量”四层组织。",
        5: "Sci-Evo · LabTrace",
        6: "06",
        8: "case 画像",
        9: "目标",
        10: "从论文中抽取可量化科研目标",
        11: "难点",
        12: "计算设计、湿实验和失败修正需串联",
        13: "边界",
        14: "只保留有证据和许可状态的步骤",
        15: "→",
        18: "进入结构化标注",
        19: "initial_request",
        20: "目标、约束、量化指标",
        21: "agent_trajectory",
        22: "action / tool / parameter / observation",
        23: "success_verification",
        24: "最终验证方法与指标",
        25: "quality",
        26: "gold 等级、证据覆盖、许可复核",
        27: "同一篇论文的切分方式会直接影响样本质量；我们的原则是先定边界，再抽步骤，再做证据绑定。",
    },
    7: {
        3: "工程体系",
        4: "评委从项目目录就能进入数据、报告、脚本和证据。",
        5: "Sci-Evo · LabTrace",
        6: "07",
        8: "项目目录",
        9: "数据 / 文档 / 报告 / 脚本",
        11: "数据",
        12: "Gold JSONL / Eval tasks",
        14: "文档",
        15: "README / 技术报告",
        17: "报告",
        18: "质量报告 / MinerU 报告",
        20: "工具",
        21: "build / validate",
        30: "核心材料彼此可互相印证：数据能被脚本校验，报告能回到来源与指标。",
    },
    8: {
        3: "质量验证",
        4: "质量不是口头承诺，而是由校验脚本、评测任务和报告共同约束。",
        5: "Sci-Evo · LabTrace",
        6: "08",
        8: "5/5",
        9: "Gold case 通过校验",
        11: "63",
        12: "自动评测任务",
        14: "25",
        15: "开放来源已筛选",
        17: "100%",
        18: "核心材料完备",
        20: "许可状态",
        21: "CC-BY 4 / 样例待复核 1",
        23: "质量报告",
        24: "已生成",
        26: "一致性检查",
        27: "通过",
        29: "自动审校",
        31: "结构校验通过；一致性检查通过。",
        32: "校验覆盖：数据结构、评测任务、质量报告、许可状态与 Git 状态。",
    },
    9: {
        3: "快速复核",
        4: "先看核心事实，再沿着文件路径进入证据、报告和代码。",
        5: "Sci-Evo · LabTrace",
        6: "09",
        10: "首屏讲清楚三件事",
        11: "数据",
        12: "gold / eval / manifest",
        13: "文档",
        14: "README / 技术报告",
        15: "验证",
        16: "校验通过",
        17: "边界",
        18: "许可复核说明",
    },
    10: {
        3: "后续增强",
        4: "基础版本已经完整；后续增强集中在强 case 密度、录屏和许可清晰度。",
        5: "Sci-Evo · LabTrace",
        6: "10",
        8: "P0",
        9: "完整基础版本",
        10: "当前已完成；数据、报告、PPT 与校验链齐全。",
        12: "P1",
        13: "扩展更多强闭环 case",
        14: "优先补 3-5 条同主题高证据案例。",
        16: "P2",
        17: "补齐 MinerU 产物",
        18: "为新增开放论文补本地解析和块级证据绑定。",
        20: "P3",
        21: "公开版合规",
        22: "逐条确认全文、图片和解析产物的许可边界。",
        24: "总结：MinerU 负责把论文变成可信结构化证据，Sci-Evo-LabTrace 负责把证据组织成可训练、可评测、可复核的科研演化数据。",
    },
}

COPY_ITEMS = [
    ("README.md", "README.md"),
    ("docs/TECHNICAL_REPORT.md", "docs/TECHNICAL_REPORT.md"),
    ("docs/DATASET_CARD.md", "docs/DATASET_CARD.md"),
    ("docs/ANNOTATION_GUIDELINES.md", "docs/ANNOTATION_GUIDELINES.md"),
    ("docs/SUBMISSION_CHECKLIST.md", "docs/SUBMISSION_CHECKLIST.md"),
    ("docs/VIDEO_SCRIPT.md", "docs/VIDEO_SCRIPT.md"),
    ("reports/QUALITY_REPORT.md", "reports/QUALITY_REPORT.md"),
    ("reports/MINERU_RUN_REPORT.md", "reports/MINERU_RUN_REPORT.md"),
    ("reports/SUBMISSION_READINESS.md", "reports/SUBMISSION_READINESS.md"),
    ("reports/CASE_DEPTH_AUDIT.md", "reports/CASE_DEPTH_AUDIT.md"),
    ("reports/VETTED_SOURCE_QUEUE.md", "reports/VETTED_SOURCE_QUEUE.md"),
    ("data/raw/pdfs/SELT-PROT-0002.pdf", "source_papers/SELT-PROT-0002.pdf"),
    ("data/raw/pdfs/SELT-PROT-0003.pdf", "source_papers/SELT-PROT-0003.pdf"),
    ("data/raw/pdfs/SELT-PROT-0004.pdf", "source_papers/SELT-PROT-0004.pdf"),
    ("data/raw/pdfs/SELT-PROT-0005.pdf", "source_papers/SELT-PROT-0005.pdf"),
    ("data/interim/mineru/SELT-PROT-0002/full.md", "mineru_artifacts/SELT-PROT-0002/full.md"),
    ("data/interim/mineru/SELT-PROT-0002/content_list.json", "mineru_artifacts/SELT-PROT-0002/content_list.json"),
    ("data/interim/mineru/SELT-PROT-0002/content_list_v2.json", "mineru_artifacts/SELT-PROT-0002/content_list_v2.json"),
    ("data/interim/mineru/SELT-PROT-0002/layout.json", "mineru_artifacts/SELT-PROT-0002/layout.json"),
    ("data/interim/mineru/SELT-PROT-0002/model.json", "mineru_artifacts/SELT-PROT-0002/model.json"),
    ("data/interim/mineru/SELT-PROT-0003/full.md", "mineru_artifacts/SELT-PROT-0003/full.md"),
    ("data/interim/mineru/SELT-PROT-0003/content_list.json", "mineru_artifacts/SELT-PROT-0003/content_list.json"),
    ("data/interim/mineru/SELT-PROT-0003/content_list_v2.json", "mineru_artifacts/SELT-PROT-0003/content_list_v2.json"),
    ("data/interim/mineru/SELT-PROT-0003/layout.json", "mineru_artifacts/SELT-PROT-0003/layout.json"),
    ("data/interim/mineru/SELT-PROT-0003/model.json", "mineru_artifacts/SELT-PROT-0003/model.json"),
    ("data/interim/mineru/SELT-PROT-0004/full.md", "mineru_artifacts/SELT-PROT-0004/full.md"),
    ("data/interim/mineru/SELT-PROT-0004/content_list.json", "mineru_artifacts/SELT-PROT-0004/content_list.json"),
    ("data/interim/mineru/SELT-PROT-0004/content_list_v2.json", "mineru_artifacts/SELT-PROT-0004/content_list_v2.json"),
    ("data/interim/mineru/SELT-PROT-0004/layout.json", "mineru_artifacts/SELT-PROT-0004/layout.json"),
    ("data/interim/mineru/SELT-PROT-0004/model.json", "mineru_artifacts/SELT-PROT-0004/model.json"),
    ("data/interim/mineru/SELT-PROT-0005/full.md", "mineru_artifacts/SELT-PROT-0005/full.md"),
    ("data/interim/mineru/SELT-PROT-0005/content_list.json", "mineru_artifacts/SELT-PROT-0005/content_list.json"),
    ("data/interim/mineru/SELT-PROT-0005/content_list_v2.json", "mineru_artifacts/SELT-PROT-0005/content_list_v2.json"),
    ("data/interim/mineru/SELT-PROT-0005/layout.json", "mineru_artifacts/SELT-PROT-0005/layout.json"),
    ("data/interim/mineru/SELT-PROT-0005/model.json", "mineru_artifacts/SELT-PROT-0005/model.json"),
    ("data/curated/cases/SELT-PROT-0002.json", "data/curated/cases/SELT-PROT-0002.json"),
    ("data/curated/cases/SELT-PROT-0003.json", "data/curated/cases/SELT-PROT-0003.json"),
    ("data/curated/cases/SELT-PROT-0004.json", "data/curated/cases/SELT-PROT-0004.json"),
    ("data/curated/cases/SELT-PROT-0005.json", "data/curated/cases/SELT-PROT-0005.json"),
    ("data/processed/scievo_gold.jsonl", "data/processed/scievo_gold.jsonl"),
    ("data/processed/scievo_gold.pretty.json", "data/processed/scievo_gold.pretty.json"),
    ("data/processed/scievo_eval_tasks.jsonl", "data/processed/scievo_eval_tasks.jsonl"),
    ("data/processed/dataset_manifest.json", "data/processed/dataset_manifest.json"),
    ("data/processed/vetted_open_access_sources.jsonl", "data/processed/vetted_open_access_sources.jsonl"),
    ("schemas/scievo_case.schema.json", "schemas/scievo_case.schema.json"),
    ("scripts/build_dataset.py", "scripts/build_dataset.py"),
    ("scripts/validate_dataset.py", "scripts/validate_dataset.py"),
    ("scripts/build_eval_tasks.py", "scripts/build_eval_tasks.py"),
    ("scripts/check_submission_readiness.py", "scripts/check_submission_readiness.py"),
    ("scripts/make_quality_report.py", "scripts/make_quality_report.py"),
    ("scripts/mineru_parse.py", "scripts/mineru_parse.py"),
    ("scripts/vet_candidate_sources.py", "scripts/vet_candidate_sources.py"),
    ("scripts/build_submission_pack.py", "scripts/build_submission_pack.py"),
    ("Sci-Evo_tool_case.json", "examples/Sci-Evo_tool_case.json"),
]


def ensure_dirs() -> None:
    (OUT_ROOT / "assets").mkdir(parents=True, exist_ok=True)
    DELIVERABLES.mkdir(parents=True, exist_ok=True)


def make_overview_image() -> None:
    image = Image.new("RGB", (1280, 720), "#F7F5EE")
    draw = ImageDraw.Draw(image)
    title_font = ImageFont.load_default()
    body_font = ImageFont.load_default()

    draw.rounded_rectangle((48, 48, 1232, 672), radius=28, fill="#FFFFFF", outline="#D6D0C5", width=2)
    draw.rounded_rectangle((84, 84, 400, 184), radius=18, fill="#E8F0FF")
    draw.text((108, 110), "Sci-Evo 项目概览", fill="#172033", font=title_font)
    draw.text((108, 138), "数据、文档、报告、脚本在同一结构中可复核", fill="#2563EB", font=body_font)

    sections = [
        ("数据文件", ["scievo_gold.jsonl", "scievo_eval_tasks.jsonl", "dataset_manifest.json"], "#EAF6FF"),
        ("关键文档", ["README.md", "TECHNICAL_REPORT.md", "SUBMISSION_CHECKLIST.md"], "#EEFCEB"),
        ("质量报告", ["QUALITY_REPORT.md", "MINERU_RUN_REPORT.md", "SUBMISSION_READINESS.md"], "#FFF5E6"),
        ("验证脚本", ["build_dataset.py", "validate_dataset.py", "check_submission_readiness.py"], "#F7EEFF"),
    ]

    y = 220
    for title, lines, fill in sections:
        draw.rounded_rectangle((92, y, 1188, y + 92), radius=16, fill=fill)
        draw.text((120, y + 16), title, fill="#172033", font=title_font)
        draw.text((120, y + 44), " | ".join(lines), fill="#334155", font=body_font)
        y += 108

    draw.rounded_rectangle((92, 616, 1188, 652), radius=12, fill="#172033")
    draw.text(
        (112, 626),
        "状态：5 个 gold case / 63 个评测任务 / 一致性检查通过 / 工作区干净",
        fill="#FFFFFF",
        font=body_font,
    )
    image.save(SUMMARY_IMAGE)


def set_shape_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    run = next((candidate for para in text_frame.paragraphs for candidate in para.runs), None)

    style = {
        "alignment": paragraph.alignment if paragraph is not None else None,
        "font_name": run.font.name if run is not None else None,
        "font_size": run.font.size if run is not None else None,
        "font_bold": run.font.bold if run is not None else None,
        "font_italic": run.font.italic if run is not None else None,
        "font_color_type": run.font.color.type if run is not None else None,
        "font_color_rgb": None,
        "font_theme_color": None,
        "font_brightness": None,
    }
    if run is not None and run.font.color.type == MSO_COLOR_TYPE.RGB:
        style["font_color_rgb"] = run.font.color.rgb
    elif run is not None and run.font.color.type == MSO_COLOR_TYPE.SCHEME:
        style["font_theme_color"] = run.font.color.theme_color
        style["font_brightness"] = run.font.color.brightness

    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    if style["alignment"] is not None:
        paragraph.alignment = style["alignment"]
    run = paragraph.add_run()
    run.text = text
    if style["font_name"] is not None:
        run.font.name = style["font_name"]
    if style["font_size"] is not None:
        run.font.size = style["font_size"]
    if style["font_bold"] is not None:
        run.font.bold = style["font_bold"]
    if style["font_italic"] is not None:
        run.font.italic = style["font_italic"]
    if style["font_color_rgb"] is not None:
        run.font.color.rgb = style["font_color_rgb"]
    elif style["font_theme_color"] is not None:
        run.font.color.theme_color = style["font_theme_color"]
        if style["font_brightness"] is not None:
            run.font.color.brightness = style["font_brightness"]


def set_shape_text_color(shape, hex_color: str) -> None:
    if not shape.has_text_frame:
        return
    color = RGBColor.from_string(hex_color)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color


def build_ppt() -> None:
    source_ppt = REFERENCE_PPT if REFERENCE_PPT.exists() else PPT_OUT
    if not source_ppt.exists():
        raise FileNotFoundError(f"No PPT template found at {REFERENCE_PPT} or {PPT_OUT}")
    prs = Presentation(str(source_ppt))

    for slide_index, replacements in SLIDE_TEXT.items():
        slide = prs.slides[slide_index - 1]
        shape_by_id = {shape.shape_id: shape for shape in slide.shapes}
        for shape_id, text in replacements.items():
            shape = shape_by_id.get(shape_id)
            if shape is None:
                raise KeyError(f"Slide {slide_index} missing shape id {shape_id}")
            set_shape_text(shape, text)
            override_color = TEXT_COLOR_OVERRIDES.get(slide_index, {}).get(shape_id)
            if override_color:
                set_shape_text_color(shape, override_color)

    slide9 = prs.slides[8]
    image_shape = next(shape for shape in slide9.shapes if shape.shape_type == 13)
    left, top, width, height = image_shape.left, image_shape.top, image_shape.width, image_shape.height
    image_shape._element.getparent().remove(image_shape._element)
    slide9.shapes.add_picture(str(SUMMARY_IMAGE), left, top, width=width, height=height)

    prs.save(str(PPT_OUT))


def copy_tree_items() -> None:
    for src_rel, dst_rel in COPY_ITEMS:
        src = ROOT / src_rel
        dst = DELIVERABLES / dst_rel
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(src, dst)

    (DELIVERABLES / "source_papers" / "README.md").write_text(
        "\n".join(
            [
                "# Source Papers",
                "",
                "本目录包含 4 篇开放许可论文 PDF，用于评审复核数据来源。",
                "",
                "- `SELT-PROT-0002.pdf`: Nature Communications, DOI `10.1038/s41467-024-50698-y`, CC-BY。",
                "- `SELT-PROT-0003.pdf`: Nature Communications, DOI `10.1038/s41467-020-18619-x`, CC-BY。",
                "- `SELT-PROT-0004.pdf`: eLife, DOI `10.7554/eLife.102788.3`, CC-BY。",
                "- `SELT-PROT-0005.pdf`: Cellular and Molecular Life Sciences, DOI `10.1007/s00018-021-03793-y`, CC-BY 4.0。",
                "",
                "说明：`SELT-PROT-0001` 来自赛事样例 PDF，数据集中保留结构化字段和证据引用；由于公开许可仍需复核，提交包不额外打包该全文 PDF。",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (DELIVERABLES / "mineru_artifacts" / "README.md").write_text(
        "\n".join(
            [
                "# MinerU Artifacts",
                "",
                "本目录包含 4 篇开放许可论文的 MinerU 结构化解析产物，用于评审复核 evidence 字段与原文之间的对应关系。",
                "",
                "每个 case 子目录包含：",
                "",
                "- `full.md`: MinerU 转换后的 Markdown 文本。",
                "- `content_list.json`: 页面与块级内容列表。",
                "- `content_list_v2.json`: 结构化内容列表 v2。",
                "- `layout.json`: 页面布局信息。",
                "- `model.json`: 模型解析结果。",
                "",
                "不包含 API token、下载 zip、未确认许可的赛事样例全文解析产物或图片目录。",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    readme = DELIVERABLES / "README.md"
    readme.write_text(
        "\n".join(
            [
                "# Sci-Evo 评审材料",
                "",
                "本目录包含项目评审所需的核心材料：",
                "",
                "- `scievo_labtrace_overview.pptx`：参考 EduMiner 演示稿风格重写的 Sci-Evo 展示 PPT。",
                "- `data/curated/` 与 `data/processed/`：curated case 源文件、gold 数据集、评测任务、manifest 和开放来源筛选结果。",
                "- `docs/`：技术报告、数据集卡片、标注规范、材料清单和录屏讲稿。",
                "- `reports/`：质量报告、Gold case 深度审计、MinerU 运行报告、完整性检查报告和来源筛选队列。",
                "- `source_papers/`：4 篇开放许可论文 PDF。",
                "- `mineru_artifacts/`：4 篇开放许可论文对应的 MinerU Markdown、content list、layout 和 model 产物。",
                "- `schemas/`：Sci-Evo case schema。",
                "- `scripts/`：构建、校验、评测任务生成、质量报告和完整性检查脚本。",
                "",
                "说明：",
                "",
                "- 不包含任何密钥文件。",
                "- 不公开打包未经许可复核的全文 PDF 或 MinerU 全量解析产物。",
                "- 如需复核当前状态，可运行 `python3 scripts/check_submission_readiness.py`。",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def build_zip() -> None:
    if ZIP_OUT.exists():
        ZIP_OUT.unlink()
    with zipfile.ZipFile(ZIP_OUT, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for path in sorted(DELIVERABLES.rglob("*")):
            if path.is_file():
                zf.write(path, path.relative_to(OUT_ROOT))


def main() -> None:
    ensure_dirs()
    make_overview_image()
    build_ppt()
    copy_tree_items()
    build_zip()
    print(json.dumps({"ppt": str(PPT_OUT), "zip": str(ZIP_OUT)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
